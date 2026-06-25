"""create_pptx — PowerPoint 슬라이드 생성 툴."""

import logging
from copy import deepcopy
from io import BytesIO
from typing import Optional, Union

from extension_modules.tools.document._common import (
    SYSTEM_TOKEN_KEYS,
    TOKEN_PATTERN,
    format_tool_response,
    load_pptx_template,
    save_to_files,
    system_token_values,
)
from langchain_core.tools import StructuredTool
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field, model_validator

log = logging.getLogger(__name__)

# Color = RGBColor (built-in theme) | MSO_THEME_COLOR (admin template inherits master theme)
ColorRef = Union[RGBColor, MSO_THEME_COLOR]

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# ── 색상 테마 (비즈니스 컨설팅풍 — Cloocus 브랜드 컬러 받으면 여기만 교체) ──
_COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)  # Deep blue — 타이틀 / 표 헤더
_COLOR_ACCENT = RGBColor(0x2E, 0x86, 0xAB)  # Teal — 부제
_COLOR_BODY = RGBColor(0x33, 0x33, 0x33)  # Dark gray — 본문 / 불릿
_COLOR_BG = RGBColor(0xFA, 0xFA, 0xFA)  # Off-white — 슬라이드 배경
_COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)  # 표 헤더 텍스트
_COLOR_TABLE_ALT = RGBColor(0xF2, 0xF2, 0xF2)  # 짝수 행 배경 (zebra)

# 폰트 크기 (Pt)
_FONT_TITLE = Pt(40)
_FONT_SUBTITLE = Pt(24)
_FONT_HEADING = Pt(28)
_FONT_BULLET = Pt(18)
_FONT_TABLE_HEADER = Pt(14)
_FONT_TABLE_BODY = Pt(12)

# 표 배치 (Title Only 레이아웃 기준)
_TABLE_LEFT = Inches(0.5)
_TABLE_TOP = Inches(1.8)
_TABLE_WIDTH = Inches(9.0)
_TABLE_HEIGHT = Inches(4.5)

# Body placeholder 인정 최소 높이. 사용자 템플릿이 banner / subtitle 라인을 type=2
# (BODY) 로 선언하는 경우 (예: '02_거버닝메시지2줄' idx=18 h=0.8in) 본문으로 오인되어
# bullets 가 좁은 배너에 쑤셔넣어지는 회귀를 차단. 표준 MS layout body 는 전부
# h ≥ 4.32in 이라 안전 마진.
_MIN_BODY_HEIGHT = Inches(2)

# Bullets textbox fallback 마진 (placeholder 못 찾았을 때). 슬라이드 너비/높이에서
# 마진/타이틀 영역만큼 빼고 사용 — widescreen 16:9 (13.33x7.5) 와 4:3 (10x7.5)
# 모두에서 컨텐츠가 좌측 ~70% 만 차지하던 회귀 수정.
_BULLET_TEXTBOX_TOP = Inches(1.5)
_BULLET_TEXTBOX_MARGIN = Inches(0.5)


def _set_color(color_obj, color: ColorRef) -> None:
    """ColorFormat 에 RGB 또는 theme color 적용."""
    if isinstance(color, RGBColor):
        color_obj.rgb = color
    else:
        color_obj.theme_color = color


def _apply_run_style(run, *, color: ColorRef, size: Pt, bold: bool = False) -> None:
    """run 단위로 폰트 색/크기/볼드 일괄 적용 (헬퍼)."""
    _set_color(run.font.color, color)
    run.font.size = size
    run.font.bold = bold


def _style_text_frame(tf, *, color: ColorRef, size: Pt, bold: bool = False) -> None:
    """text_frame 의 모든 paragraph/run 에 동일 스타일 적용."""
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            _apply_run_style(run, color=color, size=size, bold=bold)


def _set_slide_bg(slide, color: Optional[ColorRef]) -> None:
    """슬라이드 배경색 설정. ``None`` 이면 마스터 슬라이드 배경 유지."""
    if color is None:
        return
    fill = slide.background.fill
    fill.solid()
    _set_color(fill.fore_color, color)


def _iter_all_layouts(prs):
    """모든 slide master 의 모든 layout 을 순회.

    PPT 템플릿은 master 가 여러 개일 수 있고 (실제 사용자 템플릿에서 5개 master 까지 관측),
    ``prs.slide_layouts`` 는 기본 master(0번) 의 layouts 만 반환하므로 이름 매칭이 빗나간다.
    """
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            yield layout


def _layout_has_title(layout) -> bool:
    """layout 에 title placeholder (TITLE / CENTER_TITLE) 가 있는지."""
    for ph in layout.placeholders:
        try:
            t = ph.placeholder_format.type
            if t is not None and int(t) in (1, 3):  # TITLE=1, CENTER_TITLE=3
                return True
        except Exception:  # noqa: BLE001
            continue
    return False


def _layout_has_body(layout) -> bool:
    """layout 에 BODY/OBJECT placeholder 가 (title 외에) 있는지."""
    for ph in layout.placeholders:
        try:
            t = ph.placeholder_format.type
            if (
                t is not None
                and int(t) in (2, 7)  # BODY=2, OBJECT=7
                and ph.placeholder_format.idx != 0
            ):
                return True
        except Exception:  # noqa: BLE001
            continue
    return False


def _pick_layout(prs, preferred_names: tuple[str, ...], fallback_idx: int):
    """기존 시그니처 보존 (built-in 모드 호환). 이름 매칭 → title placeholder → fallback idx."""
    for name in preferred_names:
        for layout in _iter_all_layouts(prs):
            if layout.name == name:
                return layout
    for layout in _iter_all_layouts(prs):
        if _layout_has_title(layout):
            return layout
    default_master = prs.slide_masters[0]
    if 0 <= fallback_idx < len(default_master.slide_layouts):
        return default_master.slide_layouts[fallback_idx]
    return default_master.slide_layouts[0]


def _is_cover_layout(layout) -> bool:
    """cover / 표지 류 layout 식별 — body/title-only slide picker 에선 제외.

    placeholder 만 보고 판단하면 cover 에 title placeholder 추가된 케이스 (관측됨)
    가 body layout 으로 오인되어 모든 슬라이드가 cover 디자인으로 만들어진다.
    이름으로 cover 의도를 식별해 body picker 에서 제외.
    """
    name = (layout.name or "").lower()
    return (
        name.startswith("cover")
        or name.startswith("표지")
        or name.startswith("뒷표지")
        or name.startswith("back cover")
        or name.startswith("backcover")
    )


def _pick_title_layout(prs):
    """제목 슬라이드 layout — 브랜드 cover 우선.

    PPT 관례상 첫 슬라이드는 브랜드 표지 (purple cover, hero image 등) 가 적합.
    placeholder 가 없는 cover 도 ``_set_title_or_textbox`` 가 textbox 로 fallback 하므로
    visual branding 을 살리는 선택을 우선한다.

    우선순위:
    1. cover / 표지 prefix (back cover 제외) — 브랜드 표지 우선
    2. 정확 이름 — "Title Slide", "제목 슬라이드"
    3. PowerPoint 관례: Master 0 layout 0
    4. title placeholder 가 있는 첫 layout
    5. 최종 안전망: Master 0 layout 0
    """
    # 1) cover prefix (영문/한글), back cover 제외
    for layout in _iter_all_layouts(prs):
        if _is_cover_layout(layout):
            name = (layout.name or "").lower()
            if not (name.startswith("뒷표지") or name.startswith("back")):
                return layout
    # 2) 정확 이름
    for name in ("Title Slide", "제목 슬라이드"):
        for layout in _iter_all_layouts(prs):
            if layout.name == name:
                return layout
    # 3) PowerPoint 관례
    default_master = prs.slide_masters[0]
    if len(default_master.slide_layouts) > 0:
        return default_master.slide_layouts[0]
    # 4) title 있는 layout
    for layout in _iter_all_layouts(prs):
        if _layout_has_title(layout):
            return layout
    return default_master.slide_layouts[0]


def _pick_body_layout(prs):
    """본문 (bullets) 슬라이드 layout — body placeholder 우선.

    우선순위:
    1. 정확 이름 — "Title and Content", "제목 및 내용", "Contents", "내용"
    2. title + body 모두 있는 layout
    3. body 만 있는 layout (title 은 textbox 로 별도 추가)
    4. title 만 있는 layout (bullets 도 textbox 로)
    5. fallback: Master 0 layout 0
    """
    for name in ("Title and Content", "제목 및 내용", "Contents", "내용", "Content"):
        for layout in _iter_all_layouts(prs):
            if layout.name == name:
                return layout
    # cover 류는 본문에 부적합 — 명시적 제외 (cover 의 OBJECT placeholder 가 body 로
    # 잘못 매칭되는 케이스 차단)
    for layout in _iter_all_layouts(prs):
        if (
            _layout_has_title(layout)
            and _layout_has_body(layout)
            and not _is_cover_layout(layout)
        ):
            return layout
    for layout in _iter_all_layouts(prs):
        if _layout_has_body(layout) and not _is_cover_layout(layout):
            return layout
    for layout in _iter_all_layouts(prs):
        if _layout_has_title(layout) and not _is_cover_layout(layout):
            return layout
    # 최후엔 cover 라도 — 본문 전용 layout 이 아예 없는 극단적 케이스
    for layout in _iter_all_layouts(prs):
        if _layout_has_body(layout) or _layout_has_title(layout):
            return layout
    return prs.slide_masters[0].slide_layouts[0]


def _pick_title_only_layout(prs):
    """제목만 있는 layout (table 슬라이드용) — cover 제외."""
    for name in ("Title Only", "제목만"):
        for layout in _iter_all_layouts(prs):
            if layout.name == name:
                return layout
    # title 만 있고 body 없는 layout (cover 제외)
    for layout in _iter_all_layouts(prs):
        if (
            _layout_has_title(layout)
            and not _layout_has_body(layout)
            and not _is_cover_layout(layout)
        ):
            return layout
    # body 슬라이드와 동일 fallback
    return _pick_body_layout(prs)


def _clear_instance_slides(prs) -> None:
    """템플릿의 모든 인스턴스 슬라이드 제거. 마스터/레이아웃은 유지.

    python-pptx 1.0.2 에도 public 슬라이드 삭제 API 가 없어 ``_sldIdLst``
    private 인터페이스 사용 — round-trip 테스트에서 회귀 차단.
    """
    sldIdLst = prs.slides._sldIdLst
    rels_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for sldId in list(sldIdLst):
        rId = sldId.attrib[rels_ns]
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _remove_slide(prs, slide) -> None:
    """단일 슬라이드 제거 (stencil 원본 정리 등에 사용)."""
    rels_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    sldIdLst = prs.slides._sldIdLst
    target_rId = slide.part.partname
    for sldId in list(sldIdLst):
        rId = sldId.attrib[rels_ns]
        rel = prs.part.rels[rId]
        if rel.target_part.partname == target_rId:
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)
            return


def _iter_slide_text_shapes(slide):
    """슬라이드의 텍스트 가진 모든 shape (placeholder + textbox + 표 셀) yield."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        # 표 셀의 텍스트도 토큰 대상
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def _slide_has_tokens(slide) -> bool:
    """슬라이드에 `{{key}}` 패턴 1개라도 있는지 — stencil 후보 판정용."""
    for tf in _iter_slide_text_shapes(slide):
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text and "{{" in run.text:
                    return True
            # paragraph 자체 text 도 체크 (run 없이 paragraph.text 만 있는 케이스)
            if para.text and "{{" in para.text and not para.runs:
                return True
    return False


def _collect_slide_tokens(slide) -> set[str]:
    """슬라이드에서 발견된 모든 `{{key}}` 키 set."""
    keys: set[str] = set()
    for tf in _iter_slide_text_shapes(slide):
        if "{{" in tf.text:
            keys.update(TOKEN_PATTERN.findall(tf.text))
    return keys


def _collect_pptx_template_tokens(prs) -> set[str]:
    """프레젠테이션의 모든 인스턴스 슬라이드에서 발견된 토큰 키 set."""
    keys: set[str] = set()
    for slide in prs.slides:
        keys.update(_collect_slide_tokens(slide))
    return keys


def _substitute_slide_tokens(slide, tokens: dict[str, str]) -> None:
    """슬라이드의 모든 텍스트 frame 에서 `{{key}}` 를 tokens[key] 로 치환.

    run 레벨로 치환 — paragraph.text 를 새로 set 하면 run 들이 평탄화되어 폰트/색
    스타일이 깨짐. run 단위로 처리해 admin 이 디자인한 폰트/굵기/색 보존.

    누락 키는 원본 `{{key}}` 유지 (silent empty 방지).
    """
    if not tokens:
        return

    def _replace(match):
        return tokens.get(match.group(1), match.group(0))

    for tf in _iter_slide_text_shapes(slide):
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text and "{{" in run.text:
                    new_text = TOKEN_PATTERN.sub(_replace, run.text)
                    if new_text != run.text:
                        run.text = new_text


# OOXML relationships namespace — r:id/r:embed/r:link 속성에 사용.
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_R_ID_ATTRS = (
    f"{{{_R_NS}}}id",
    f"{{{_R_NS}}}embed",
    f"{{{_R_NS}}}link",
)
_LAYOUT_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)


def _duplicate_slide(prs, source_slide):
    """source_slide 의 디자인을 같은 layout 으로 복제한 새 슬라이드 반환.

    python-pptx 는 native ``slide.copy()`` 가 없어 layout 기반 새 슬라이드 생성 +
    shape XML deepcopy 패턴 사용. layout 으로부터 상속되는 placeholder text 는
    중복으로 들어가지 않도록 새 슬라이드의 빈 placeholder 들을 먼저 제거.

    Source 의 ``part.rels`` 도 새 슬라이드에 복제 + 매핑된 rId 로 XML 속성 rewrite —
    image (`r:embed`) / 차트 / 하이퍼링크 (`r:id`) / OLE embed 등이 박힌 stencil 이
    복제 후에도 정상 렌더되도록. slideLayout rel 은 ``add_slide`` 가 이미 설정한
    것이라 skip (중복 등록 시 PowerPoint 가 파일 오픈 거부).
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    sp_tree = new_slide.shapes._spTree
    # layout 에서 상속받은 placeholder shape 제거 (id 충돌 + 중복 방지)
    for shp in list(new_slide.shapes):
        sp_tree.remove(shp.element)

    # source 의 rels (slideLayout 제외) 복제 — 새 rId 가 자동 할당됨
    rid_map: dict[str, str] = {}
    for src_rid, rel in source_slide.part.rels.items():
        if rel.reltype == _LAYOUT_RELTYPE:
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True
            )
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        if new_rid != src_rid:
            rid_map[src_rid] = new_rid

    # source shape XML deepcopy + r:id/r:embed/r:link 속성을 새 rId 로 rewrite
    for shape in source_slide.shapes:
        cloned = deepcopy(shape.element)
        if rid_map:
            for elem in cloned.iter():
                for attr in _R_ID_ATTRS:
                    old_rid = elem.get(attr)
                    if old_rid is not None and old_rid in rid_map:
                        elem.set(attr, rid_map[old_rid])
        sp_tree.append(cloned)
    return new_slide


def _discover_pptx_template_token_keys() -> set[str]:
    """현재 활성 PPT 템플릿을 로드해 발견된 모든 토큰 키 set 반환.

    템플릿 미설정 / 파싱 실패 / IO 에러는 모두 empty set 으로 swallow.
    """
    try:
        prs = load_pptx_template()
        if prs is None:
            return set()
        return _collect_pptx_template_tokens(prs)
    except Exception as e:  # noqa: BLE001
        log.warning("Failed to discover PPT template tokens: %s", e)
        return set()


def _add_styled_table(
    slide,
    table_data: list[list[str]],
    *,
    header_bg: ColorRef,
    header_text: ColorRef,
    body_text: ColorRef,
    odd_row_bg: Optional[ColorRef] = None,
    even_row_bg: Optional[ColorRef] = None,
) -> None:
    """슬라이드에 표 추가. 색상 인자로 받아 built-in 테마 / 템플릿 테마 모두 지원."""
    rows = len(table_data)
    cols = len(table_data[0]) if rows else 0
    if rows == 0 or cols == 0:
        return

    shape = slide.shapes.add_table(
        rows, cols, _TABLE_LEFT, _TABLE_TOP, _TABLE_WIDTH, _TABLE_HEIGHT
    )
    table = shape.table

    for r_idx, row_data in enumerate(table_data):
        for c_idx, value in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = "" if value is None else str(value)

            if r_idx == 0:
                cell.fill.solid()
                _set_color(cell.fill.fore_color, header_bg)
                _style_text_frame(
                    cell.text_frame,
                    color=header_text,
                    size=_FONT_TABLE_HEADER,
                    bold=True,
                )
            else:
                row_bg = odd_row_bg if r_idx % 2 == 0 else even_row_bg
                if row_bg is not None:
                    cell.fill.solid()
                    _set_color(cell.fill.fore_color, row_bg)
                _style_text_frame(
                    cell.text_frame, color=body_text, size=_FONT_TABLE_BODY
                )


class PPTSlide(BaseModel):
    title: str = Field(..., max_length=200)
    bullets: list[str] = Field(default_factory=list, max_length=20)
    table: Optional[list[list[str]]] = Field(
        None,
        max_length=30,  # 행 상한 — 슬라이드 한 장에 들어가는 시각적 한계 + 메모리 가드
        description=(
            "표 데이터 (첫 행=헤더). 지정 시 해당 슬라이드는 Title Only 레이아웃으로 "
            "렌더링되며 bullets 는 자동 무시된다. 행 ≤ 12, 열 ≤ 8 권장."
        ),
    )
    speaker_notes: Optional[str] = Field(None, max_length=2000)

    @model_validator(mode="after")
    def _normalize_content_choice(self):
        # bullets 와 table 동시 지정 시 table 우선 — bullets 는 silent drop.
        # 명시적 reject 가 아니라 normalize 로 처리해 LLM 이 둘 다 보내도 동작.
        if self.bullets and self.table:
            self.bullets = []
        return self


class PPTContent(BaseModel):
    filename: str = Field(..., description="확장자 제외, 한글/영문 가능")
    title: str = Field(..., description="첫 슬라이드(타이틀) 제목")
    subtitle: Optional[str] = Field(None, description="첫 슬라이드 부제")
    template_tokens: Optional[dict[str, str]] = Field(
        None,
        description=(
            "관리자 PPT 템플릿의 표지 (1번째 토큰 슬라이드) `{{key}}` 자리표시자를 "
            "런타임에 치환할 값. 예: `{'title': '2026 Q1 매출 보고서'}`. "
            "`date`/`author` 등 시스템 토큰은 자동 채워지므로 LLM 이 컨텐츠 토큰만 "
            "신경 쓰면 된다. 키 누락 시 `{{key}}` 가 그대로 남는다."
        ),
    )
    stencil_slides: Optional[list[dict[str, str]]] = Field(
        None,
        max_length=50,
        description=(
            "관리자 PPT 템플릿의 2번째 이후 토큰 슬라이드(stencil)를 N번 복제하면서 "
            "각 항목 token 값으로 치환. 예: 본문 섹션 헤더 디자인 1장 + LLM 이 보낸 "
            "`[{'title': '1장', 'content': '...'}, {'title': '2장', 'content': '...'}]` "
            "→ 같은 디자인의 슬라이드 2장 생성. 단일 텍스트 토큰용 — bullets/표가 "
            "필요한 본문은 아래 `slides` 사용. 최대 50장 (deepcopy DoS 차단)."
        ),
    )
    slides: list[PPTSlide] = Field(
        default_factory=list,
        max_length=100,
        description=(
            "본문 컨텐츠 슬라이드 (title + bullets/표). 빌더가 body layout 으로 새 "
            "슬라이드 생성. template_tokens / stencil_slides 만 쓰는 시나리오에선 "
            "비워도 됨 (빈 발표자료는 model_validator 가 차단)."
        ),
    )

    @model_validator(mode="after")
    def _require_some_content(self):
        # 표지 토큰 + stencil + slides 중 하나는 있어야 의미 있는 PPT 생성.
        # 비어있으면 빈 표지 슬라이드 1장만 나오는 가짜 성공 차단.
        if (
            not self.slides
            and not self.stencil_slides
            and not (self.template_tokens or self.subtitle)
        ):
            raise ValueError(
                "At least one of `slides`, `stencil_slides`, or `template_tokens` "
                "must be provided"
            )
        return self


# ── Builder ──────────────────────────────────────────────────────────────────


def _build_pptx(
    content: PPTContent,
    *,
    system_tokens: Optional[dict[str, str]] = None,
) -> BytesIO:
    """PPTContent → in-memory .pptx BytesIO.

    템플릿 미설정 시: 비즈니스 컨설팅풍 hard-coded 테마.
    템플릿 설정 시: 마스터 슬라이드 / 테마 색상 / 폰트를 inherit 하고
        텍스트는 ``MSO_THEME_COLOR.ACCENT_1`` 등 theme color enum 으로 칠해
        admin 이 업로드한 brand 색상이 자동 적용된다.

    템플릿 인스턴스 슬라이드 처리:
    - 토큰 `{{key}}` 박힌 슬라이드는 보존 + 치환
    - 첫 토큰 슬라이드 = 표지 (single use, ``template_tokens`` 로 치환)
    - 2번째 이후 토큰 슬라이드 = 본문 stencil (``stencil_slides`` 의 엔트리만큼 복제)
    - 토큰 없는 인스턴스 슬라이드 = 샘플로 간주, 제거
    - 그 뒤에 ``slides`` (PPTSlide list) 를 layout 기반으로 추가 (기존 동작)
    """
    template_prs = load_pptx_template()
    template_used = template_prs is not None
    prs = template_prs if template_used else Presentation()

    cover_slide_processed = False
    stencil_templates: list = []  # 마지막 cleanup 단계에서 참조 — 빈 list 가 기본
    if template_used:
        # 1) BEFORE any modification: classify slides as token-bearing vs not.
        # 치환 후엔 토큰이 사라지므로 분류는 반드시 먼저.
        all_slides = list(prs.slides)
        token_slides = [s for s in all_slides if _slide_has_tokens(s)]
        non_token_slides = [s for s in all_slides if not _slide_has_tokens(s)]
        cover_slide = token_slides[0] if token_slides else None
        stencil_templates = token_slides[1:]

        # 2) 토큰 없는 인스턴스 슬라이드 (샘플/장식) 제거 — stencil 복제하기 전에 정리.
        for slide in non_token_slides:
            _remove_slide(prs, slide)

        # 3) 표지(cover) in-place 치환. LLM 이 template_tokens.title 안 보내도
        # content.title 을 fallback 으로 주입해 표지가 비지 않게 한다.
        if cover_slide is not None:
            merged_cover_tokens: dict[str, str] = {
                **(system_tokens or {}),
                "title": content.title,
                **(content.template_tokens or {}),
            }
            if content.subtitle is not None:
                merged_cover_tokens.setdefault("subtitle", content.subtitle)
            _substitute_slide_tokens(cover_slide, merged_cover_tokens)
            cover_slide_processed = True

        # 4) Stencil 복제 — v1 은 단일 stencil 만 사용 (다중 stencil 은 추후 라벨링 도입 시).
        # 원본 stencil 제거는 모든 add_slide 끝난 뒤 마지막에 — python-pptx 의
        # `_next_partname_idx = len(slides)+1` 가 슬라이드 중간 제거 시 다음 add_slide
        # 에서 이미 사용된 partname 을 재할당해 두 슬라이드가 같은 part 가리키는 손상이
        # 발생함 (직접 관측).
        if stencil_templates and content.stencil_slides:
            stencil = stencil_templates[0]
            if len(stencil_templates) > 1:
                log.info(
                    "PPT template has %d stencil slides — v1 uses only the first",
                    len(stencil_templates),
                )
            for entry in content.stencil_slides:
                merged_entry = {**(system_tokens or {}), **entry}
                clone = _duplicate_slide(prs, stencil)
                _substitute_slide_tokens(clone, merged_entry)
        # 원본 stencil 제거는 LLM body slides add_slide 다음에 — 본 함수 끝에서 수행.

    # 색상: 템플릿 사용 시 theme color enum (마스터 색상 inherit),
    # 미사용 시 hard-coded RGB.
    if template_used:
        title_color: ColorRef = MSO_THEME_COLOR.ACCENT_1
        subtitle_color: ColorRef = MSO_THEME_COLOR.ACCENT_2
        heading_color: ColorRef = MSO_THEME_COLOR.ACCENT_1
        body_text_color: ColorRef = MSO_THEME_COLOR.TEXT_1
        slide_bg: Optional[ColorRef] = None  # 마스터 배경 유지
        table_header_bg: ColorRef = MSO_THEME_COLOR.ACCENT_1
        table_odd_row_bg: Optional[ColorRef] = None  # 마스터 zebra 유지
        table_even_row_bg: Optional[ColorRef] = None
    else:
        title_color = _COLOR_PRIMARY
        subtitle_color = _COLOR_ACCENT
        heading_color = _COLOR_PRIMARY
        body_text_color = _COLOR_BODY
        slide_bg = _COLOR_BG
        table_header_bg = _COLOR_PRIMARY
        table_odd_row_bg = _COLOR_TABLE_ALT
        table_even_row_bg = _COLOR_BG

    title_layout = _pick_title_layout(prs)
    body_layout = _pick_body_layout(prs)
    title_only_layout = _pick_title_only_layout(prs)

    # 템플릿 모드에선 placeholder 텍스트에 explicit run formatting 을 적용하지 않는다.
    # 이유: python-pptx 의 placeholder 는 run<-paragraph<-layout<-master 순서로 상속하는데
    # _style_text_frame() 이 run.font.{color,size,bold} 를 setting 하면 master 의 폰트/색/사이즈가
    # 모두 덮어씌워져서 사용자가 업로드한 템플릿 디자인이 사라진다.
    def _style_placeholder(tf, *, color, size, bold=False):
        if template_used:
            return  # placeholder 가 master 로부터 상속하도록 둠
        _style_text_frame(tf, color=color, size=size, bold=bold)

    def _find_title_placeholder(slide):
        """슬라이드의 title 자리 placeholder.

        우선순위: TITLE(1) / CENTER_TITLE(3) → SUBTITLE(4) → BODY(2) / OBJECT(7).

        사용자 cover 레이아웃이 단일 OBJECT placeholder 로 메인 텍스트 영역을 디자인한
        케이스 (관측됨) 에서도 title 이 그 placeholder 로 들어가야 시각적으로 올바름.
        body 슬라이드는 _find_body_placeholder 가 exclude_ph 로 title 차지한 걸 건너뛰어
        bullets 가 충돌 안 함.
        """
        subtitle_ph = None
        body_ph = None
        for ph in slide.placeholders:
            try:
                t = ph.placeholder_format.type
                if t is None:
                    continue
                ti = int(t)
                if ti in (1, 3):  # TITLE / CENTER_TITLE — 최우선
                    return ph
                if ti == 4 and subtitle_ph is None:  # SUBTITLE
                    subtitle_ph = ph
                elif ti in (2, 7) and body_ph is None:  # BODY / OBJECT
                    body_ph = ph
            except Exception:  # noqa: BLE001
                continue
        return subtitle_ph or body_ph

    def _set_title_or_textbox(slide, text, *, color, size, bold):
        """title placeholder (또는 SUBTITLE fallback) 가 있으면 그걸 사용, 없으면 textbox.

        Returns the placeholder used (or None if textbox) so caller 가 bullets 단계에서
        같은 placeholder 를 다시 쓰지 않도록 exclude 할 수 있다.
        """
        title_ph = _find_title_placeholder(slide)
        if title_ph is not None:
            title_ph.text = text
            _style_placeholder(title_ph.text_frame, color=color, size=size, bold=bold)
            return title_ph
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tb.text_frame.text = text
        # textbox 는 placeholder 가 아니라 master 상속 없음 → 항상 explicit 스타일링
        _style_text_frame(tb.text_frame, color=color, size=size, bold=bold)
        return None

    def _find_body_placeholder(slide, exclude=None):
        """slide.placeholders 중 BODY/OBJECT type 이고 본문 컨텐츠 영역으로 보이는
        placeholder 중 **면적이 가장 큰 것** 선택. 없으면 None (caller textbox fallback).

        사용자 템플릿의 일부 layout 은 banner / subtitle 라인을 ``type=2 (BODY)``
        로 선언한다 (관측됨: '02_거버닝메시지2줄' 의 idx=18 type=2 top=0.7in h=0.8in
        — 본문이 아니라 거버닝 메시지 2번째 줄 배너). 본문으로 오인하면 모든 bullets
        가 0.8in 짜리 좁은 배너에 쑤셔넣어진다 — 정확히 사용자 보고 사례.

        Heuristic: 본문은 최소 ``_MIN_BODY_HEIGHT`` 이상이어야 한다. 표준 MS 레이아웃의
        body 는 전부 height ≥ 4.32in (Comparison content 가 가장 낮음) 라 안전 마진.
        layout 에 단일 BODY 만 있는 경우에도 banner 면 제외 → textbox fallback.

        ``exclude.placeholder_format.idx`` 와 동일한 idx 는 건너뜀 (title 중복 방지).
        """
        exclude_idx = exclude.placeholder_format.idx if exclude is not None else None
        candidates = []
        for ph in slide.placeholders:
            try:
                if exclude_idx is not None and ph.placeholder_format.idx == exclude_idx:
                    continue
                t = ph.placeholder_format.type
                if t is None or int(t) not in (2, 7):  # BODY=2, OBJECT=7
                    continue
                if ph.placeholder_format.idx == 0:  # 0=title 제외
                    continue
                # width / height 가 None 일 수 있음 (inherited from layout) — 0 처리
                w = ph.width or 0
                h = ph.height or 0
                if h < _MIN_BODY_HEIGHT:  # banner-style — 본문 아님
                    continue
                area = int(w) * int(h)
                candidates.append((area, ph))
            except Exception:  # noqa: BLE001
                continue
        if not candidates:
            return None
        # 면적 최대 — 동일 면적이면 iteration 순서 유지 (sort 안정성)
        candidates.sort(key=lambda x: x[0], reverse=True)
        return candidates[0][1]

    def _set_bullets_or_textbox(slide, bullets, *, color, size, exclude_ph=None):
        """body placeholder 가 있으면 거기에, 없으면 textbox 로 bullets 렌더.

        textbox fallback 은 슬라이드 크기에서 좌우 마진과 상단 타이틀 영역만 빼고
        나머지 전체 — widescreen 템플릿에서 컨텐츠가 좌측만 차지하는 회귀 차단.
        """
        body_ph = _find_body_placeholder(slide, exclude=exclude_ph)
        if body_ph is not None:
            tf = body_ph.text_frame
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
            _style_placeholder(tf, color=color, size=size)
            return
        tb_width = prs.slide_width - 2 * _BULLET_TEXTBOX_MARGIN
        tb_height = prs.slide_height - _BULLET_TEXTBOX_TOP - _BULLET_TEXTBOX_MARGIN
        tb = slide.shapes.add_textbox(
            _BULLET_TEXTBOX_MARGIN, _BULLET_TEXTBOX_TOP, tb_width, tb_height
        )
        tf = tb.text_frame
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
        _style_text_frame(tf, color=color, size=size)

    # 첫 슬라이드: title layout — 단, 템플릿의 토큰 표지가 이미 채워졌으면 skip.
    # (그 표지가 이미 첫 슬라이드 역할을 함)
    if not cover_slide_processed:
        title_slide = prs.slides.add_slide(title_layout)
        _set_slide_bg(title_slide, slide_bg)

        title_used_ph = _set_title_or_textbox(
            title_slide, content.title, color=title_color, size=_FONT_TITLE, bold=True
        )
    else:
        title_slide = None
        title_used_ph = None

    if content.subtitle and not cover_slide_processed:
        # title 이 차지한 placeholder idx 와 다른 첫 텍스트 placeholder 찾기.
        # python-pptx 는 placeholders[idx] 가 매 호출마다 새 wrapper 를 반환할 수 있어
        # `is` 비교 대신 idx 로 식별.
        title_used_idx = (
            title_used_ph.placeholder_format.idx if title_used_ph is not None else None
        )
        subtitle_ph = None
        for ph in title_slide.placeholders:
            if (
                title_used_idx is not None
                and ph.placeholder_format.idx == title_used_idx
            ):
                continue
            try:
                t = ph.placeholder_format.type
                if t is not None and int(t) in (1, 2, 3, 4, 7):
                    subtitle_ph = ph
                    break
            except Exception:  # noqa: BLE001
                continue
        if subtitle_ph is not None:
            subtitle_ph.text = content.subtitle
            _style_placeholder(
                subtitle_ph.text_frame, color=subtitle_color, size=_FONT_SUBTITLE
            )

    for slide_def in content.slides:
        layout = title_only_layout if slide_def.table else body_layout
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, slide_bg)

        used_title_ph = _set_title_or_textbox(
            slide, slide_def.title, color=heading_color, size=_FONT_HEADING, bold=True
        )

        if slide_def.table:
            _add_styled_table(
                slide,
                slide_def.table,
                header_bg=table_header_bg,
                header_text=_COLOR_WHITE,
                body_text=body_text_color,
                odd_row_bg=table_odd_row_bg,
                even_row_bg=table_even_row_bg,
            )
        elif slide_def.bullets:
            _set_bullets_or_textbox(
                slide,
                slide_def.bullets,
                color=body_text_color,
                size=_FONT_BULLET,
                exclude_ph=used_title_ph,
            )

        if slide_def.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_def.speaker_notes

    # 원본 stencil 슬라이드는 가장 마지막에 제거 — add_slide 가 끝난 후라야
    # partname 재할당 충돌 없음. stencil 사용 여부와 무관하게 정리.
    if template_used:
        for stencil_slide in stencil_templates:
            _remove_slide(prs, stencil_slide)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


_TOOL_DESCRIPTION = (
    "Generate a PowerPoint presentation (.pptx) file. "
    "발표 자료/슬라이드/프레젠테이션을 .pptx 파일로 생성한다. "
    "Use when the user asks for: 'PPT', 'PowerPoint', 'presentation', 'slides', "
    "'slide deck', 'deck', '슬라이드', '발표자료', '제안서 슬라이드' etc. "
    "Each slide takes a title plus EITHER `bullets` (list of bullet strings) OR "
    "`table` (2D list, first row = header) — use `table` when the slide content is "
    "tabular (비교표/지표/일정표 등). Optional `speaker_notes` per slide."
)


def _build_dynamic_description() -> str:
    """기본 description 에 현재 PPT 템플릿에서 발견된 컨텐츠 토큰 키 안내 덧붙임.

    토큰이 1번째 슬라이드(표지)와 2번째 슬라이드(stencil)에 다르게 분포할 수 있어
    여기선 단순히 발견된 모든 컨텐츠 토큰을 나열. 시스템 자동 토큰은 제외.
    """
    base = _TOOL_DESCRIPTION
    discovered = _discover_pptx_template_token_keys()
    content_keys = sorted(discovered - SYSTEM_TOKEN_KEYS)
    if content_keys:
        keys_csv = ", ".join(content_keys)
        base += (
            " IMPORTANT — the admin's PPT template uses these content placeholders: "
            f"{keys_csv}. Fill them via `template_tokens` (cover, single use) or "
            "`stencil_slides` (body section, list of dicts for N-time duplication). "
            "Date/author tokens are auto-filled by the system."
        )
    return base


def make_create_pptx(user_id: str) -> StructuredTool:
    """user_id에 바인딩된 create_pptx 툴을 생성."""

    def _create_pptx(**kwargs) -> str:
        log.info(
            "create_pptx called: user_id=%s, template_tokens=%r, "
            "stencil_slides=%d, slides=%d",
            user_id,
            kwargs.get("template_tokens"),
            len(kwargs.get("stencil_slides") or []),
            len(kwargs.get("slides") or []),
        )
        content = PPTContent(**kwargs)
        system_tokens = system_token_values(user_id)
        buf = _build_pptx(content, system_tokens=system_tokens)
        result = save_to_files(
            user_id=user_id,
            filename=f"{content.filename}.pptx",
            buffer=buf,
            mime=PPTX_MIME,
        )
        return format_tool_response(result)

    return StructuredTool.from_function(
        func=_create_pptx,
        name="create_pptx",
        description=_build_dynamic_description(),
        args_schema=PPTContent,
    )
