"""Document generation tool common helpers."""

import logging
import os
import re
import time
import uuid
from datetime import datetime, timezone
from io import BytesIO
from typing import Any, Optional
from zoneinfo import ZoneInfo

from open_webui.config import (
    DOCUMENT_TEMPLATE_DOCX,
    DOCUMENT_TEMPLATE_PPTX,
    DOCUMENT_TEMPLATE_XLSX,
    WEBUI_URL,
)
from open_webui.env import SRC_LOG_LEVELS
from open_webui.models.files import FileForm, Files
from open_webui.models.users import Users
from open_webui.routers.files import generate_signed_url
from open_webui.storage.provider import Storage
from pydantic import BaseModel, Field

# 표지/스텐실 슬라이드/시트의 `{{key}}` 자리표시자 — pptx_tool/xlsx_tool 공통.
TOKEN_PATTERN = re.compile(r"\{\{(\w+)\}\}")

# 시스템이 자동 주입하는 토큰 키 — LLM 이 안 보내도 채워짐. 컨텐츠 토큰과 구분해
# tool description 에서 빼주기 위한 set.
SYSTEM_TOKEN_KEYS: frozenset[str] = frozenset(
    {
        "date",
        "today",
        "now",
        "datetime",
        "year",
        "month",
        "day",
        "author",
        "user",
        "user_name",
        "name",
        "user_email",
        "email",
    }
)

log = logging.getLogger(__name__)
log.setLevel(SRC_LOG_LEVELS.get("MAIN", "INFO"))

# 파일명 타임스탬프 표준 시간대 — 환경변수 TZ 가 있으면 그걸로, 없으면 KST.
_FILENAME_TZ = ZoneInfo(os.environ.get("TZ", "Asia/Seoul"))

# 토큰 치환용 시간대 (TZ env 우선) — pptx/xlsx 양쪽에서 사용.
_TOKEN_TZ = _FILENAME_TZ


def system_token_values(user_id: Optional[str]) -> dict[str, str]:
    """시스템 자동 주입 토큰 값 — 사용자/날짜 컨텍스트.

    LLM 이 ``template_tokens`` 를 안 보내도 ``{{date}}``, ``{{author}}`` 류는
    여기서 자동으로 채워진다. LLM 이 같은 키를 보내면 그쪽이 우선 (override).
    pptx_tool / xlsx_tool 공통 사용.
    """
    now = datetime.now(_TOKEN_TZ)
    values: dict[str, str] = {
        "date": now.strftime("%Y-%m-%d"),
        "today": now.strftime("%Y-%m-%d"),
        "now": now.strftime("%Y-%m-%d %H:%M"),
        "datetime": now.strftime("%Y-%m-%d %H:%M"),
        "year": str(now.year),
        "month": f"{now.month:02d}",
        "day": f"{now.day:02d}",
    }
    if user_id:
        try:
            user = Users.get_user_by_id(user_id)
            if user:
                name = (user.name or user.email or "").strip()
                email = (user.email or "").strip()
                values.update(
                    {
                        "author": name,
                        "user": name,
                        "user_name": name,
                        "name": name,
                        "user_email": email,
                        "email": email,
                    }
                )
        except Exception as e:  # noqa: BLE001 — DB 못 읽어도 빌드 계속
            log.warning("Failed to load user for system tokens (%s): %s", user_id, e)
    return values


class DocumentToolOutput(BaseModel):
    """모든 문서 생성 툴의 공통 응답."""

    file_id: str = Field(..., description="Files 테이블 PK (UUID)")
    filename: str = Field(..., description="확장자 포함 최종 파일명")
    download_url: str = Field(
        ..., description="다운로드 URL: /api/v1/files/{id}/content"
    )
    size_bytes: int = Field(..., ge=0, description="파일 크기 (바이트)")


_FILENAME_SAFE_RE = re.compile(r"[^\w\-_.]", re.UNICODE)
_MAX_FILENAME_LEN = 100


def sanitize_filename(name: str) -> str:
    """파일명에서 경로 분리자/특수문자 제거. 한글/영문/숫자/_-./ 만 유지.

    빈 문자열이거나 정제 후 비면 'document' 로 폴백. 길이 100자 제한.
    """
    if not name:
        return "document"
    safe = _FILENAME_SAFE_RE.sub("_", name)
    return safe[:_MAX_FILENAME_LEN]


def _build_canonical_filename(filename: str, user_email: str) -> str:
    """{베이스명}_{이메일}_{yymmddhhmmss}.{확장자} 형태의 표준 파일명 생성.

    Storage(blob 또는 local) 키와 DB filename(다운로드 시 노출되는 이름)이
    동일하게 이 형태로 저장된다. timezone 은 KST 기본 (TZ env 으로 override).
    """
    safe_input = sanitize_filename(filename)
    if "." in safe_input:
        base, ext = safe_input.rsplit(".", 1)
    else:
        base, ext = safe_input, ""

    safe_email = sanitize_filename(user_email or "anonymous") or "anonymous"
    ts = datetime.now(timezone.utc).astimezone(_FILENAME_TZ).strftime("%y%m%d%H%M%S")

    canonical = f"{base}_{safe_email}_{ts}"
    if ext:
        canonical = f"{canonical}.{ext}"
    # 모든 구성요소가 이미 sanitized 라 regex 재실행은 no-op — 길이 캡만 적용.
    return canonical[:_MAX_FILENAME_LEN]


def save_to_files(
    user_id: str,
    filename: str,
    buffer: BytesIO,
    mime: str,
) -> DocumentToolOutput:
    """Storage 업로드 + Files 테이블 등록을 한 번에 처리.

    Storage 키와 DB filename 은 모두 `{베이스명}_{이메일}_{yymmddhhmmss}.확장자`
    형태로 통일. 다운로드 시 사용자에게 동일한 파일명으로 노출됨.
    Files DB 저장 실패 시 RuntimeError.
    """
    user = Users.get_user_by_id(user_id) if user_id else None
    user_email = (user.email if user else None) or user_id or "anonymous"

    canonical_name = _build_canonical_filename(filename, user_email)
    file_id = str(uuid.uuid4())

    buffer.seek(0)
    contents, file_path = Storage.upload_file(buffer, canonical_name)

    form = FileForm(
        id=file_id,
        filename=canonical_name,
        path=file_path,
        meta={"content_type": mime, "size": len(contents), "source": "document_tool"},
    )
    file_model = Files.insert_new_file(user_id, form)
    if file_model is None:
        raise RuntimeError(f"Files 테이블 저장 실패: {canonical_name}")

    # Signed URL — 브라우저가 인증 헤더 없이 클릭만으로 다운로드 가능 + 1시간 만료.
    # WEBUI_URL 이 절대 base 라 vite dev proxy(5173) 우회. schedule_worker.py 와
    # 동일 패턴 (chart_image upload).
    base_url = (WEBUI_URL.value or "http://localhost:8080").rstrip("/")
    signed_url = generate_signed_url(base_url, file_model.id)

    return DocumentToolOutput(
        file_id=file_model.id,
        filename=canonical_name,
        download_url=signed_url,
        size_bytes=len(contents),
    )


# ────────────────────────────────────────────────────────────────────────────
# Template loading (PR1 — admin-uploaded master templates)
# ────────────────────────────────────────────────────────────────────────────

# Module-level cache: kind → (uploaded_at, raw_bytes). Single slot per kind
# is enough since each config holds exactly one current template. Storage
# round-trip on every _build_*() would block the event loop and add 100-200ms
# per LLM tool invocation on Azure Blob — we cache the BYTES (not the parsed
# object) so each call still gets a fresh Presentation/Document/Workbook
# instance (those are not safe to share across builds).
_TEMPLATE_CACHE: dict[str, tuple[int, bytes]] = {}

# Storage IOError retry window — Azure Blob has rare BlobNotFound right after
# write due to eventual consistency. Single retry is enough; second failure
# falls back to built-in theme.
_STORAGE_RETRY_DELAY_S = 0.2

_KIND_TO_CONFIG = {
    "pptx": DOCUMENT_TEMPLATE_PPTX,
    "docx": DOCUMENT_TEMPLATE_DOCX,
    "xlsx": DOCUMENT_TEMPLATE_XLSX,
}


def _get_template_bytes(kind: str) -> Optional[bytes]:
    """Fetch raw template bytes for ``kind`` from config + Storage.

    Returns None when no template is configured OR Storage fetch fails
    repeatedly OR the config dict is malformed. Callers MUST handle None
    by falling back to the built-in theme (i.e. ``Presentation()``).

    Cache key is ``uploaded_at`` — when admin replaces the template,
    ``uploaded_at`` changes and we re-fetch automatically.
    """
    persistent_config = _KIND_TO_CONFIG.get(kind)
    if persistent_config is None:
        return None

    cfg = persistent_config.value
    if not isinstance(cfg, dict):
        return None

    file_path = cfg.get("file_path")
    uploaded_at = cfg.get("uploaded_at", 0)
    if not file_path:
        return None

    # Cache hit?
    cached = _TEMPLATE_CACHE.get(kind)
    if cached and cached[0] == uploaded_at:
        return cached[1]

    # Storage fetch with single retry (Azure eventual consistency)
    # ``Storage.get_file`` 의 모든 프로바이더 구현은 로컬 파일 경로(str)를 반환한다
    # (LocalStorageProvider 는 원본 경로 그대로, S3/GCS/Azure 는 UPLOAD_DIR 로 다운로드
    # 후 그 경로). 따라서 반환값을 그대로 BytesIO 에 넘기면 TypeError → silent fallback.
    # glossary._load_file_bytes() 와 동일한 read-then-bytes 패턴.
    for attempt in (1, 2):
        try:
            local_path = Storage.get_file(file_path)
            with open(local_path, "rb") as f:
                raw = f.read()
            _TEMPLATE_CACHE[kind] = (uploaded_at, raw)
            return raw
        except Exception as e:  # noqa: BLE001 — provider-specific exceptions
            if attempt == 1:
                log.warning(
                    "Storage fetch failed for %s template (attempt 1): %s",
                    kind,
                    e,
                )
                time.sleep(_STORAGE_RETRY_DELAY_S)
                continue
            log.error(
                "Storage fetch failed for %s template (attempt 2, giving up): %s",
                kind,
                e,
            )
            return None
    return None


def load_pptx_template() -> Optional[Any]:
    """Load admin-configured PPT template as a ``Presentation`` instance.

    Returns None when no template is configured or template is corrupt —
    caller falls back to ``Presentation()`` (built-in theme).
    """
    raw = _get_template_bytes("pptx")
    if raw is None:
        return None
    try:
        from pptx import Presentation

        return Presentation(BytesIO(raw))
    except Exception as e:  # noqa: BLE001 — BadZipFile, PackageNotFoundError, ...
        log.warning("PPTX template parse failed, falling back: %s", e)
        return None


def load_docx_template() -> Optional[Any]:
    """Load admin-configured Word template as a ``Document`` instance."""
    raw = _get_template_bytes("docx")
    if raw is None:
        return None
    try:
        from docx import Document

        return Document(BytesIO(raw))
    except Exception as e:  # noqa: BLE001
        log.warning("DOCX template parse failed, falling back: %s", e)
        return None


def load_xlsx_template() -> Optional[Any]:
    """Load admin-configured Excel template as a ``Workbook`` instance."""
    raw = _get_template_bytes("xlsx")
    if raw is None:
        return None
    try:
        from openpyxl import load_workbook

        return load_workbook(BytesIO(raw))
    except Exception as e:  # noqa: BLE001
        log.warning("XLSX template parse failed, falling back: %s", e)
        return None


# UnifiedAgent 가 ToolMessage 에서 우리 응답을 식별해 LLM 응답 전에 사용자에게
# 직접 스트림 yield 하기 위한 안정 마커. 변경 시 unified_agent 의 detection 로직도
# 함께 업데이트.
DOCUMENT_TOOL_MARKER = "[document_tool] 파일 생성 완료"


def format_tool_response(output: DocumentToolOutput) -> str:
    """LLM 이 응답에 그대로 인용할 수 있는 마크다운 한 줄로 변환.

    첫 줄의 DOCUMENT_TOOL_MARKER 는 UnifiedAgent 가 ToolMessage 를 식별해 LLM
    응답 전에 사용자에게 직접 스트림 yield 하기 위한 안정 토큰이다 (image_tool
    의 `![Generated Image]` 패턴과 동일).

    LLM 의 hallucination ('파일 생성 도구 없음') 차단은 final answer 시스템
    프롬프트의 'TOOL EXECUTION SUCCESSFUL' 섹션이 담당한다 — 여기서 별도
    행동 지시 텍스트를 두면 중복.
    """
    size_kb = output.size_bytes / 1024
    return (
        f"{DOCUMENT_TOOL_MARKER}: [{output.filename}]({output.download_url}) "
        f"({size_kb:,.1f} KB)"
    )
