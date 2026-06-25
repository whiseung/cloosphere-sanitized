"""Document template inheritance tests — PR1 of admin-document-templates.

Covers the new ``load_*_template()`` helpers in ``_common.py`` and the
builder branches in ``pptx_tool.py`` / ``docx_tool.py`` / ``xlsx_tool.py``.

Strategy: stub the ``DOCUMENT_TEMPLATE_*`` PersistentConfig objects' ``.value``
and the ``Storage.get_file`` call, so we exercise the full flow without
needing a real DB or storage provider.
"""

import os
import tempfile
from contextlib import contextmanager
from io import BytesIO
from unittest.mock import patch

import pytest
from docx import Document as DocxOpen
from extension_modules.tools.document import _common
from extension_modules.tools.document._common import (
    _TEMPLATE_CACHE,
    _get_template_bytes,
    load_pptx_template,
)
from extension_modules.tools.document.docx_tool import (
    DocxBlock,
    DocxContent,
    _build_docx,
)
from extension_modules.tools.document.pptx_tool import (
    PPTContent,
    PPTSlide,
    _build_pptx,
)
from extension_modules.tools.document.xlsx_tool import (
    XlsxContent,
    XlsxSheet,
    _build_xlsx,
)
from openpyxl import Workbook
from openpyxl import load_workbook as xlsx_open
from pptx import Presentation


@pytest.fixture(autouse=True)
def _clear_template_cache():
    """Cache lives at module level; reset between tests to avoid bleed."""
    _TEMPLATE_CACHE.clear()
    yield
    _TEMPLATE_CACHE.clear()


def _empty_pptx_bytes() -> bytes:
    prs = Presentation()
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _empty_docx_bytes() -> bytes:
    doc = DocxOpen()
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _empty_xlsx_bytes() -> bytes:
    wb = Workbook()
    wb.active.title = "TemplateSheet"
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _set_template_config(kind: str, *, uploaded_at: int = 1700000000) -> None:
    """Set the module-level PersistentConfig to a stubbed dict."""
    config = _common._KIND_TO_CONFIG[kind]
    config.value = {
        "version": 1,
        "file_path": f"document-templates/{kind}/test.{kind}",
        "original_filename": f"test.{kind}",
        "uploaded_at": uploaded_at,
        "uploaded_by": "u-admin",
    }


def _clear_template_config(kind: str) -> None:
    _common._KIND_TO_CONFIG[kind].value = {}


@contextmanager
def _mock_storage_returns_path(raw: bytes):
    """Storage.get_file 의 실제 시그니처(=로컬 파일 경로 반환)를 정확히 재현.

    이전엔 ``return_value=raw`` (bytes) 로 mock 했으나, 실제 구현은 path(str) 반환
    이라 production 에서만 TypeError 가 나는 false-pass 버그가 있었음.
    """
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".tpl")
    try:
        tmp.write(raw)
        tmp.close()
        with patch.object(
            _common.Storage, "get_file", return_value=tmp.name
        ) as mock_get:
            yield mock_get
    finally:
        if os.path.exists(tmp.name):
            os.unlink(tmp.name)


# ── Cache & retry tests ──────────────────────────────────────────────────────


class TestTemplateBytesCache:
    def test_no_config_returns_none(self):
        _clear_template_config("pptx")
        assert _get_template_bytes("pptx") is None

    def test_cache_hit_avoids_storage_call(self):
        raw = _empty_pptx_bytes()
        _set_template_config("pptx", uploaded_at=42)
        with _mock_storage_returns_path(raw) as mock_get:
            _get_template_bytes("pptx")  # miss
            _get_template_bytes("pptx")  # hit
            _get_template_bytes("pptx")  # hit
            assert mock_get.call_count == 1

    def test_cache_invalidated_on_uploaded_at_change(self):
        raw = _empty_pptx_bytes()
        _set_template_config("pptx", uploaded_at=1)
        with _mock_storage_returns_path(raw) as mock_get:
            _get_template_bytes("pptx")  # miss → fetch
            _set_template_config("pptx", uploaded_at=2)  # admin replaced
            _get_template_bytes("pptx")  # miss → fetch again
            assert mock_get.call_count == 2

    def test_storage_error_retried_once_then_returns_none(self):
        _set_template_config("pptx")
        with patch.object(
            _common.Storage,
            "get_file",
            side_effect=IOError("blob not found"),
        ) as mock_get:
            with patch.object(_common.time, "sleep", return_value=None):
                result = _get_template_bytes("pptx")
            assert result is None
            assert mock_get.call_count == 2

    def test_storage_recovers_on_second_attempt(self):
        raw = _empty_pptx_bytes()
        _set_template_config("pptx")
        # 첫 시도 실패 → 두 번째 시도엔 실제 파일 경로 반환
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".tpl")
        try:
            tmp.write(raw)
            tmp.close()
            with patch.object(
                _common.Storage,
                "get_file",
                side_effect=[IOError("blob not found"), tmp.name],
            ) as mock_get:
                with patch.object(_common.time, "sleep", return_value=None):
                    result = _get_template_bytes("pptx")
                assert result == raw
                assert mock_get.call_count == 2
        finally:
            if os.path.exists(tmp.name):
                os.unlink(tmp.name)

    def test_corrupt_template_falls_back_to_none(self):
        """load_*_template() catches library exceptions, returns None."""
        _set_template_config("pptx")
        with _mock_storage_returns_path(b"NOT A ZIP FILE"):
            assert load_pptx_template() is None


# ── Regression: builders preserve existing behavior when no template ─────────


class TestBuildersWithoutTemplate:
    """When no template is configured, builders must behave exactly as before."""

    def setup_method(self):
        for k in ("pptx", "docx", "xlsx"):
            _clear_template_config(k)

    def test_pptx_no_template_round_trip(self):
        content = PPTContent(
            filename="d", title="T", slides=[PPTSlide(title="S", bullets=["x"])]
        )
        buf = _build_pptx(content)
        buf.seek(0)
        prs = Presentation(buf)
        assert len(prs.slides) == 2  # title + 1
        assert prs.slides[0].shapes.title.text == "T"

    def test_docx_no_template_round_trip(self):
        content = DocxContent(
            filename="d",
            title="Doc Title",
            blocks=[DocxBlock(type="paragraph", text="body")],
        )
        buf = _build_docx(content)
        buf.seek(0)
        doc = DocxOpen(buf)
        texts = [p.text for p in doc.paragraphs]
        assert "Doc Title" in texts and "body" in texts

    def test_xlsx_no_template_round_trip(self):
        content = XlsxContent(
            filename="d",
            sheets=[
                XlsxSheet(name="S1", columns=["A", "B"], rows=[["1", "2"]]),
            ],
        )
        buf = _build_xlsx(content)
        buf.seek(0)
        wb = xlsx_open(buf)
        assert wb.sheetnames == ["S1"]


# ── Template mode behavior ────────────────────────────────────────────────────


class TestBuildersWithTemplate:
    def test_pptx_template_used_clears_instance_slides_and_inherits_layouts(self):
        """Template's layouts/masters preserved, instance slides cleared."""
        raw = _empty_pptx_bytes()
        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(
                filename="d",
                title="With Template",
                slides=[PPTSlide(title="S", bullets=["a"])],
            )
            buf = _build_pptx(content)
        buf.seek(0)
        prs = Presentation(buf)
        # Empty template starts with 0 slides; we added 2 (title + 1 content)
        assert len(prs.slides) == 2
        # Layouts are inherited from template (≥11 in MS standard masters)
        assert len(prs.slide_layouts) >= 6

    def test_pptx_template_mode_skips_explicit_run_formatting(self):
        """Regression: 템플릿 모드에선 placeholder run 의 explicit rPr 가 비어 있어야
        master 의 폰트/사이즈/색이 상속된다. 이전엔 _style_text_frame 이 무조건
        호출되어 master 스타일이 모두 덮어씌워졌다 (사용자 보고 — 'PPT 템플릿 스타일 안먹음').
        """
        raw = _empty_pptx_bytes()
        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(
                filename="d",
                title="Hello",
                slides=[PPTSlide(title="Slide", bullets=["one"])],
            )
            buf = _build_pptx(content)
        buf.seek(0)
        prs = Presentation(buf)

        # Title slide: title placeholder 의 run 이 explicit color/size/bold 없이 inherit
        title_run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
        assert title_run.font.size is None, "size 가 set 되면 master 폰트 사이즈 무시됨"
        assert title_run.font.bold is None, "bold 가 set 되면 master bold 설정 무시됨"
        # color 는 _NoneColor 또는 inherited — explicit rgb/theme 셋팅이 없어야 함
        assert title_run.font.color.type is None, (
            "color.type 이 set 되면 master color scheme 무시됨"
        )

        # Body slide: heading placeholder 도 동일
        body_run = prs.slides[1].shapes.title.text_frame.paragraphs[0].runs[0]
        assert body_run.font.size is None
        assert body_run.font.bold is None
        assert body_run.font.color.type is None

    def test_pptx_no_template_applies_explicit_styling(self):
        """반대 케이스: 템플릿 미사용 시엔 built-in 테마 (RGB) 가 explicit 적용돼야 한다."""
        for k in ("pptx",):
            _clear_template_config(k)
        content = PPTContent(
            filename="d", title="T", slides=[PPTSlide(title="S", bullets=["a"])]
        )
        buf = _build_pptx(content)
        buf.seek(0)
        prs = Presentation(buf)
        title_run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
        assert title_run.font.size is not None
        assert title_run.font.bold is True
        assert title_run.font.color.type is not None

    def test_pptx_template_with_existing_slides_clears_them(self):
        """Templates that ship with placeholder slides — those get removed."""
        # Build a template with 3 pre-existing slides
        tpl = Presentation()
        for _ in range(3):
            tpl.slides.add_slide(tpl.slide_layouts[0])
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(filename="d", title="T", slides=[PPTSlide(title="S")])
            buf = _build_pptx(content)
        buf.seek(0)
        prs = Presentation(buf)
        # Should only have OUR slides (title + 1), not 3+2
        assert len(prs.slides) == 2

    def test_docx_template_appends_after_existing_content(self):
        """DOCX mode preserves template content (form headers etc.)."""
        # Build a template with a heading "FORM HEADER"
        tpl = DocxOpen()
        tpl.add_heading("FORM HEADER", level=1)
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("docx")
        with _mock_storage_returns_path(raw):
            content = DocxContent(
                filename="d",
                title="New Body",
                blocks=[DocxBlock(type="paragraph", text="body text")],
            )
            buf = _build_docx(content)
        buf.seek(0)
        doc = DocxOpen(buf)
        texts = [p.text for p in doc.paragraphs]
        # Template content preserved + our content appended
        assert "FORM HEADER" in texts
        assert "New Body" in texts
        assert "body text" in texts

    def test_xlsx_template_preserves_existing_sheets(self):
        """XLSX mode preserves template sheets (formulas, named ranges etc.)."""
        raw = _empty_xlsx_bytes()  # has one sheet 'TemplateSheet'
        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            content = XlsxContent(
                filename="d",
                sheets=[XlsxSheet(name="NewData", columns=["A"], rows=[["1"]])],
            )
            buf = _build_xlsx(content)
        buf.seek(0)
        wb = xlsx_open(buf)
        # Template sheet + our sheet both present
        assert "TemplateSheet" in wb.sheetnames
        assert "NewData" in wb.sheetnames

    def test_xlsx_template_tokens_substituted_on_cover_sheet(self):
        """템플릿 표지 시트의 {{key}} 가 template_tokens 값으로 치환된다."""
        # 표지 시트가 있는 템플릿 생성 — A1, B2, A3 에 토큰 박아둠
        tpl = Workbook()
        cover = tpl.active
        cover.title = "Cover"
        cover["A1"] = "{{title}}"
        cover["B2"] = "작성일: {{date}}"
        cover["A3"] = "작성자: {{author}}"
        cover["A4"] = "정적 텍스트"  # 치환 대상 아님
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            content = XlsxContent(
                filename="d",
                template_tokens={
                    "title": "2026 Q1 매출 보고서",
                    "date": "2026-05-20",
                    "author": "홍길동",
                },
                sheets=[XlsxSheet(name="Data", columns=["A"], rows=[["x"]])],
            )
            buf = _build_xlsx(content)

        buf.seek(0)
        wb = xlsx_open(buf)
        cover_ws = wb["Cover"]
        assert cover_ws["A1"].value == "2026 Q1 매출 보고서"
        assert cover_ws["B2"].value == "작성일: 2026-05-20"
        assert cover_ws["A3"].value == "작성자: 홍길동"
        assert cover_ws["A4"].value == "정적 텍스트"
        # 데이터 시트는 그대로
        assert "Data" in wb.sheetnames
        assert wb["Data"]["A4"].value == "x"  # 데이터 행 (row 4~)

    def test_xlsx_template_missing_token_keeps_placeholder(self):
        """누락된 키는 원본 {{key}} 가 그대로 남는다 (silent empty 방지)."""
        tpl = Workbook()
        cover = tpl.active
        cover.title = "Cover"
        cover["A1"] = "{{title}} - {{missing}}"
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            content = XlsxContent(
                filename="d",
                template_tokens={"title": "Report"},
                sheets=[XlsxSheet(name="Data", columns=["A"], rows=[["x"]])],
            )
            buf = _build_xlsx(content)
        buf.seek(0)
        wb = xlsx_open(buf)
        assert wb["Cover"]["A1"].value == "Report - {{missing}}"

    def test_xlsx_template_tokens_do_not_touch_llm_sheets(self):
        """LLM 이 만든 데이터 시트의 합법적인 {{...}} 텍스트는 치환되지 않는다."""
        tpl = Workbook()
        tpl.active.title = "Cover"
        tpl.active["A1"] = "{{title}}"
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            content = XlsxContent(
                filename="d",
                template_tokens={"title": "Report"},
                sheets=[
                    XlsxSheet(
                        name="Data",
                        columns=["설명"],
                        rows=[["{{title}} 자리표시자 예시"]],
                    )
                ],
            )
            buf = _build_xlsx(content)
        buf.seek(0)
        wb = xlsx_open(buf)
        assert wb["Cover"]["A1"].value == "Report"
        # LLM 시트 본문에 박힌 토큰 텍스트는 그대로 유지
        assert wb["Data"]["A4"].value == "{{title}} 자리표시자 예시"

    def test_xlsx_formula_injection_escaped_in_data_rows(self):
        """회귀: LLM 이 보낸 데이터 셀이 Excel 수식으로 평가되지 않도록 prefix 이스케이프.

        =, +, -, @, Tab, CR 로 시작하는 셀은 Excel 이 수식으로 해석 — `=WEBSERVICE(...)`
        같은 외부 호출이 파일 여는 사용자 컨텍스트에서 실행됨. 보안 감사 VULN-001.
        """
        for k in ("pptx", "docx", "xlsx"):
            _clear_template_config(k)
        content = XlsxContent(
            filename="d",
            sheets=[
                XlsxSheet(
                    name="S",
                    columns=["A", "B", "C", "D"],
                    rows=[
                        [
                            '=WEBSERVICE("https://atk")',  # 수식 시작
                            "+1+2",
                            "@SUM(A1)",
                            "normal",  # 정상 데이터 — 변형되지 않아야
                        ],
                    ],
                ),
            ],
        )
        buf = _build_xlsx(content)
        buf.seek(0)
        wb = xlsx_open(buf)
        ws = wb["S"]
        # 데이터 row 는 row 4 (DATA_START_ROW)
        assert ws["A4"].value == '\'=WEBSERVICE("https://atk")'  # ' prefix
        assert ws["B4"].value == "'+1+2"
        assert ws["C4"].value == "'@SUM(A1)"
        assert ws["D4"].value == "normal"  # 정상값은 그대로

    def test_xlsx_formula_injection_escaped_in_template_tokens(self):
        """토큰 치환 결과도 LLM 출처라 동일하게 이스케이프돼야 함."""
        tpl = Workbook()
        tpl.active.title = "Cover"
        tpl.active["A1"] = "{{payload}}"
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            content = XlsxContent(
                filename="d",
                template_tokens={"payload": '=HYPERLINK("https://atk","OK")'},
                sheets=[XlsxSheet(name="D", columns=["A"], rows=[["x"]])],
            )
            buf = _build_xlsx(content)
        buf.seek(0)
        wb = xlsx_open(buf)
        assert wb["Cover"]["A1"].value == '\'=HYPERLINK("https://atk","OK")'

    def test_pptx_stencil_slides_max_length_enforced(self):
        """DoS 방지: stencil_slides 가 50개 초과면 Pydantic validation 차단."""
        import pytest as _pytest

        with _pytest.raises(Exception) as exc_info:
            PPTContent(
                filename="d",
                title="T",
                stencil_slides=[{"title": f"s{i}"} for i in range(51)],
                slides=[PPTSlide(title="S", bullets=["x"])],
            )
        # Pydantic ValidationError — "List should have at most 50 items"
        assert "50" in str(exc_info.value) or "max" in str(exc_info.value).lower()

    def test_xlsx_template_tokens_ignored_without_template(self):
        """템플릿 미사용 시 template_tokens 는 무시 (에러 없음)."""
        for k in ("pptx", "docx", "xlsx"):
            _clear_template_config(k)
        content = XlsxContent(
            filename="d",
            template_tokens={"title": "Report"},
            sheets=[XlsxSheet(name="Data", columns=["A"], rows=[["x"]])],
        )
        buf = _build_xlsx(content)  # must not raise
        buf.seek(0)
        wb = xlsx_open(buf)
        assert wb.sheetnames == ["Data"]

    def test_xlsx_system_tokens_auto_fill_when_llm_silent(self):
        """LLM 이 template_tokens 안 보내도 시스템 토큰 (date/author) 자동 채움.

        실제 사용자 보고 시나리오: '{{date}}/{{author}}/{{title}} 다 그대로 나옴'
        — LLM 이 template_tokens 를 안 보내는 게 흔한 케이스. 시스템 토큰만으로도
        2/3 가 자동 채워져야 한다 (title 은 LLM 만 알 수 있어 unfilled).
        """
        tpl = Workbook()
        tpl.active.title = "Cover"
        tpl.active["A1"] = "{{title}}"
        tpl.active["A2"] = "{{author}}"
        tpl.active["A3"] = "{{date}}"
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            content = XlsxContent(
                filename="d",
                template_tokens=None,  # LLM 안 보냄
                sheets=[XlsxSheet(name="Data", columns=["A"], rows=[["x"]])],
            )
            system_tokens = {"date": "2026-05-20", "author": "홍길동"}
            buf = _build_xlsx(content, system_tokens=system_tokens)
        buf.seek(0)
        wb = xlsx_open(buf)
        cover = wb["Cover"]
        assert cover["A1"].value == "{{title}}"  # LLM 만 채울 수 있음
        assert cover["A2"].value == "홍길동"  # 시스템 자동
        assert cover["A3"].value == "2026-05-20"  # 시스템 자동

    def test_xlsx_llm_token_overrides_system_token(self):
        """LLM 이 보낸 값이 시스템 값보다 우선 (specificity)."""
        tpl = Workbook()
        tpl.active.title = "Cover"
        tpl.active["A1"] = "{{date}}"
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            content = XlsxContent(
                filename="d",
                template_tokens={"date": "2030-01-01"},  # LLM 의 명시값
                sheets=[XlsxSheet(name="Data", columns=["A"], rows=[["x"]])],
            )
            system_tokens = {"date": "2026-05-20"}
            buf = _build_xlsx(content, system_tokens=system_tokens)
        buf.seek(0)
        wb = xlsx_open(buf)
        assert wb["Cover"]["A1"].value == "2030-01-01"

    def test_docx_cover_token_substitution(self):
        """DOCX 템플릿 표지의 {{key}} 가 치환되고 content.title 중복 paragraph 는 skip."""
        tpl = DocxOpen()
        tpl.add_paragraph("{{title}}", style="Title")
        tpl.add_paragraph("작성자: {{author}}")
        tpl.add_paragraph("작성일: {{date}}")
        tpl.add_page_break()
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("docx")
        with _mock_storage_returns_path(raw):
            content = DocxContent(
                filename="d",
                title="UNUSED — cover token has its own title",
                template_tokens={"title": "Report Title"},
                blocks=[DocxBlock(type="paragraph", text="body text")],
            )
            buf = _build_docx(
                content, system_tokens={"author": "Hong", "date": "2026-05-20"}
            )
        buf.seek(0)
        out = DocxOpen(buf)
        texts = [p.text for p in out.paragraphs]
        # 표지 토큰 치환됨
        assert "Report Title" in texts
        assert "작성자: Hong" in texts
        assert "작성일: 2026-05-20" in texts
        # content.title (UNUSED) 의 중복 Title paragraph 없음 — cover 가 이미 title
        assert "UNUSED — cover token has its own title" not in texts
        # 본문 추가됨
        assert "body text" in texts

    def test_docx_missing_token_keeps_placeholder(self):
        """누락 토큰은 원본 {{key}} 유지 (admin 오타 진단)."""
        tpl = DocxOpen()
        tpl.add_paragraph("{{title}} - {{missing}}")
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("docx")
        with _mock_storage_returns_path(raw):
            content = DocxContent(
                filename="d",
                title="Report",
                template_tokens={"title": "Report"},
                blocks=[DocxBlock(type="paragraph", text="x")],
            )
            buf = _build_docx(content)
        buf.seek(0)
        out = DocxOpen(buf)
        texts = [p.text for p in out.paragraphs]
        assert "Report - {{missing}}" in texts

    def test_docx_tokens_in_textbox_substituted(self):
        """회귀: Word 표지 디자인이 textbox/shape 안에 있는 경우에도 토큰 치환.

        python-docx 의 `doc.paragraphs` 는 flow content (body, table cell) 만
        보고 textbox 안 paragraph 는 안 보여줌. 실제 사용자 보고: `{{title}}` 가
        DrawingML/VML textbox 안에 있어 치환 실패. _iter_textbox_paragraphs 가
        XML 트리에서 직접 찾아야 동작.
        """
        from docx.oxml.ns import qn
        from lxml import etree

        # python-docx 로 textbox 만들기 어려움 — XML 직접 주입
        tpl = DocxOpen()
        tpl.add_paragraph("body before")  # 일반 body paragraph
        # body 의 마지막 paragraph 다음에 textbox 삽입 (DrawingML)
        body_elem = tpl.element.body
        textbox_xml = """
        <w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
             xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
             xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
          <w:r>
            <w:drawing>
              <wp:inline distT="0" distB="0" distL="0" distR="0">
                <wp:extent cx="5760720" cy="914400"/>
                <wp:docPr id="1" name="Cover Title Box"/>
                <a:graphic>
                  <a:graphicData uri="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
                    <wps:wsp>
                      <wps:cNvSpPr/>
                      <wps:spPr/>
                      <wps:txbx>
                        <w:txbxContent>
                          <w:p><w:r><w:t>{{title}}</w:t></w:r></w:p>
                        </w:txbxContent>
                      </wps:txbx>
                      <wps:bodyPr/>
                    </wps:wsp>
                  </a:graphicData>
                </a:graphic>
              </wp:inline>
            </w:drawing>
          </w:r>
        </w:p>
        """
        # sectPr 직전에 삽입
        sectpr = body_elem.find(qn("w:sectPr"))
        body_elem.insert(
            list(body_elem).index(sectpr), etree.fromstring(textbox_xml.strip())
        )

        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("docx")
        with _mock_storage_returns_path(raw):
            content = DocxContent(
                filename="d",
                title="Cover Title",
                template_tokens={"title": "Cover Title"},
                blocks=[DocxBlock(type="paragraph", text="body content")],
            )
            buf = _build_docx(content)
        buf.seek(0)
        out = DocxOpen(buf)

        # Textbox 안에서 {{title}} 가 치환됐는지 확인 — XML 레벨 점검
        from docx.oxml.ns import qn as _qn

        W_TXBX = _qn("w:txbxContent")
        W_T = _qn("w:t")
        textbox_texts = []
        for txbx in out.element.body.iter(W_TXBX):
            text = "".join(t.text or "" for t in txbx.iter(W_T))
            textbox_texts.append(text)
        assert "Cover Title" in textbox_texts, (
            f"textbox token 치환 실패 — {{title}} 그대로 남음 (textbox_texts={textbox_texts})"
        )

    def test_docx_template_without_tokens_preserves_existing_behavior(self):
        """토큰 없는 템플릿: content.title 의 Title paragraph 가 그대로 append (회귀 차단)."""
        raw = _empty_docx_bytes()  # no tokens
        _set_template_config("docx")
        with _mock_storage_returns_path(raw):
            content = DocxContent(
                filename="d",
                title="My Title",
                blocks=[DocxBlock(type="paragraph", text="body")],
            )
            buf = _build_docx(content)
        buf.seek(0)
        out = DocxOpen(buf)
        texts = [p.text for p in out.paragraphs]
        assert "My Title" in texts  # 기존 동작 유지
        assert "body" in texts

    def test_pptx_cover_token_substitution(self):
        """PPT 템플릿의 표지 슬라이드 (첫 토큰 슬라이드) 의 {{key}} 가 치환된다."""
        from pptx.util import Inches

        tpl = Presentation()
        tpl.slide_width = Inches(13.33)
        tpl.slide_height = Inches(7.5)
        # 첫 슬라이드에 토큰 textbox 박기
        cover = tpl.slides.add_slide(tpl.slide_layouts[5])  # Title Only
        tb = cover.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tb.text_frame.text = "{{title}}"
        tb2 = cover.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.5))
        tb2.text_frame.text = "{{author}} / {{date}}"
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(
                filename="d",
                title="UNUSED",
                template_tokens={"title": "Report Title"},
                slides=[PPTSlide(title="Body", bullets=["x"])],
            )
            buf = _build_pptx(
                content, system_tokens={"author": "Hong", "date": "2026-05-20"}
            )
        buf.seek(0)
        out = Presentation(buf)
        # Cover (slide 0) — 토큰 치환됨
        cover_texts = [
            s.text_frame.text for s in out.slides[0].shapes if s.has_text_frame
        ]
        assert "Report Title" in cover_texts
        assert "Hong / 2026-05-20" in cover_texts

    def test_pptx_stencil_duplication_creates_n_clones(self):
        """2번째 토큰 슬라이드(stencil)가 stencil_slides 엔트리 수만큼 복제된다.

        실제 사용자 시나리오 — partname 충돌 방지 회귀 포함 (python-pptx
        _next_partname_idx 가 슬라이드 중간 제거 시 같은 partname 재할당하던 버그).
        """
        from pptx.util import Inches

        tpl = Presentation()
        tpl.slide_width = Inches(13.33)
        tpl.slide_height = Inches(7.5)
        # 첫 슬라이드 = cover token
        cover = tpl.slides.add_slide(tpl.slide_layouts[5])
        tb = cover.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tb.text_frame.text = "{{title}}"
        # 두번째 슬라이드 = stencil token
        stencil = tpl.slides.add_slide(tpl.slide_layouts[5])
        tb1 = stencil.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tb1.text_frame.text = "{{title}}"
        tb2 = stencil.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        tb2.text_frame.text = "{{content}}"
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(
                filename="d",
                title="UNUSED",
                template_tokens={"title": "Report"},
                stencil_slides=[
                    {"title": "S1", "content": "C1"},
                    {"title": "S2", "content": "C2"},
                    {"title": "S3", "content": "C3"},
                ],
                slides=[PPTSlide(title="Body", bullets=["x"])],
            )
            buf = _build_pptx(content, system_tokens=None)
        buf.seek(0)
        out = Presentation(buf)
        # cover + 3 stencil clones + 1 LLM body = 5 slides
        assert len(out.slides) == 5
        # 모든 partname 이 unique 해야 (회귀: slide5.xml 가 두 번 할당되던 버그)
        partnames = [s.part.partname for s in out.slides]
        assert len(set(partnames)) == 5, f"partname collision: {partnames}"

        # 각 stencil 클론은 다른 title/content
        sten1_texts = [
            s.text_frame.text for s in out.slides[1].shapes if s.has_text_frame
        ]
        sten2_texts = [
            s.text_frame.text for s in out.slides[2].shapes if s.has_text_frame
        ]
        sten3_texts = [
            s.text_frame.text for s in out.slides[3].shapes if s.has_text_frame
        ]
        assert "S1" in sten1_texts and "C1" in sten1_texts
        assert "S2" in sten2_texts and "C2" in sten2_texts
        assert "S3" in sten3_texts and "C3" in sten3_texts
        # 원본 stencil 은 제거됨 — 같은 토큰 텍스트가 살아있지 않아야
        for s in out.slides:
            for shape in s.shapes:
                if shape.has_text_frame:
                    assert "{{content}}" not in shape.text_frame.text
                    assert "{{title}}" not in shape.text_frame.text

    def test_pptx_stencil_duplicate_preserves_image_rels(self):
        """회귀: stencil 슬라이드에 이미지가 박혀있으면 복제 시 r:embed → image
        part rel 도 복제돼야 한다. shape XML 만 deepcopy 하면 cloned 의 rId 가
        존재하지 않는 참조가 되어 PowerPoint 가 'broken image' 표시 또는 파일 오픈
        거부. 운영 사고 1순위 시나리오 (admin 이 로고/일러스트 박은 stencil)."""
        from io import BytesIO as _BIO

        from pptx.util import Inches

        tpl = Presentation()
        tpl.slide_width = Inches(13.33)
        tpl.slide_height = Inches(7.5)

        # 1px PNG 바이트 (이미지 추가용)
        png_bytes = bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
            "890000000d49444154789c63f80f0000010101006e9e5e8e0000000049454e44ae426082"
        )

        # 첫 슬라이드 = cover (단순 token)
        cover = tpl.slides.add_slide(tpl.slide_layouts[5])
        cover.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(1)
        ).text_frame.text = "{{title}}"

        # 두번째 슬라이드 = stencil with IMAGE — 회귀 핵심
        stencil = tpl.slides.add_slide(tpl.slide_layouts[5])
        stencil.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(1)
        ).text_frame.text = "{{title}}"
        stencil.shapes.add_picture(
            _BIO(png_bytes), Inches(2), Inches(3), Inches(2), Inches(2)
        )

        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(
                filename="d",
                title="UNUSED",
                template_tokens={"title": "Cover"},
                stencil_slides=[{"title": "S1"}, {"title": "S2"}],
                slides=[PPTSlide(title="Body", bullets=["x"])],
            )
            buf = _build_pptx(content)
        buf.seek(0)
        out = Presentation(buf)

        # 각 clone 의 picture shape 가 정상 image part 를 참조하는지 검증.
        # python-pptx 의 picture.image 접근이 broken rId 면 KeyError 발생.
        from pptx.oxml.ns import qn

        BLIP_QN = qn("a:blip")
        R_EMBED_QN = (
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
        )
        for si in (1, 2):  # 두 stencil clone
            slide = out.slides[si]
            blips_with_embed = [
                blip for blip in slide.element.iter(BLIP_QN) if blip.get(R_EMBED_QN)
            ]
            assert blips_with_embed, f"Slide {si}: 이미지 blip 없음"
            for blip in blips_with_embed:
                rid = blip.get(R_EMBED_QN)
                # rId 가 slide.part.rels 에 존재해야 함 — 없으면 broken image
                assert rid in slide.part.rels, (
                    f"Slide {si}: blip r:embed={rid} 가 part.rels 에 없음 — "
                    f"이미지 stencil 복제 회귀"
                )
                # target_part 가 image content_type 인지도 확인
                target = slide.part.rels[rid].target_part
                assert "image" in target.content_type, (
                    f"Slide {si}: rId={rid} target 이 image 아님 ({target.content_type})"
                )

    def test_pptx_template_without_tokens_falls_back_to_layout(self):
        """토큰 없는 템플릿: 기존 동작 그대로 — 모든 인스턴스 슬라이드 제거 후
        layout 기반 생성 (회귀 차단)."""
        raw = _empty_pptx_bytes()
        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(
                filename="d",
                title="T",
                slides=[PPTSlide(title="S", bullets=["a"])],
            )
            buf = _build_pptx(content)
        buf.seek(0)
        out = Presentation(buf)
        assert len(out.slides) == 2  # title + 1 body

    def test_xlsx_dynamic_description_lists_content_tokens(self):
        """_build_dynamic_description 이 발견된 컨텐츠 토큰만 안내한다."""
        from extension_modules.tools.document.xlsx_tool import (
            _build_dynamic_description,
        )

        tpl = Workbook()
        tpl.active.title = "Cover"
        tpl.active["A1"] = "{{title}}"
        tpl.active["A2"] = "{{summary}}"  # 컨텐츠 토큰
        tpl.active["A3"] = "{{date}}"  # 시스템 토큰 — 안내에서 제외돼야
        tpl.active["A4"] = "{{author}}"  # 시스템 토큰
        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("xlsx")
        with _mock_storage_returns_path(raw):
            desc = _build_dynamic_description()

        # 컨텐츠 토큰은 안내에 포함
        assert "title" in desc
        assert "summary" in desc
        # 시스템 토큰은 LLM 이 채울 필요 없으니 IMPORTANT 안내에서 제외
        # ("auto-filled by the system" 문구 자체가 나타나는지로 검증)
        assert "auto-filled" in desc

    def test_pptx_banner_body_placeholder_falls_back_to_textbox(self):
        """회귀: 템플릿 layout 의 BODY type placeholder 가 banner 라인일 때
        (예: '거버닝 메시지 2줄' 의 2번째 줄 높이 0.8in) bullets 가 그 좁은
        배너에 쑤셔넣어지지 않고 textbox 로 본문 영역에 떨어진다.

        실제 사용자 보고: '모든메시지가 거버닝에 들어가네 하단에 컨텐츠영역에
        들어가야하는데'.
        """
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        tpl = Presentation()
        # 16:9 widescreen — 사용자 템플릿과 동일
        tpl.slide_width = Inches(13.33)
        tpl.slide_height = Inches(7.5)
        # 첫 layout 의 placeholder 들을 banner 만으로 교체.
        # layout XML 을 직접 조작해 SUBTITLE 라인 1개 + BODY banner 2개 구성
        # (사용자 템플릿 '02_거버닝메시지2줄' 의 placeholder shape 와 동일 구조).
        layout = tpl.slide_layouts[1]  # Title and Content — 본문 placeholder 제거
        # 기존 placeholder 전부 제거
        spTree = layout.shapes._spTree
        for sp in list(spTree.iter(qn("p:sp"))):
            spTree.remove(sp)

        # SUBTITLE (type=subTitle, idx=17) — 거버닝 메시지 1번째 줄
        # BODY (type=body, idx=15) — 작은 사이드 placeholder
        # BODY (type=body, idx=18) — 거버닝 메시지 2번째 줄 (banner, h=0.8in)
        layout_xml = """
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="100" name="Subtitle"/>
            <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
            <p:nvPr><p:ph type="subTitle" idx="17"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr>
            <a:xfrm>
              <a:off x="731520" y="91440"/>
              <a:ext cx="10241280" cy="457200"/>
            </a:xfrm>
          </p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
        body_small_xml = """
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="101" name="SmallBody"/>
            <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
            <p:nvPr><p:ph type="body" idx="15" sz="quarter"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr>
            <a:xfrm>
              <a:off x="9692640" y="182880"/>
              <a:ext cx="2011680" cy="182880"/>
            </a:xfrm>
          </p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
        body_banner_xml = """
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="102" name="BannerLine2"/>
            <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
            <p:nvPr><p:ph type="body" idx="18"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr>
            <a:xfrm>
              <a:off x="731520" y="640080"/>
              <a:ext cx="10881360" cy="731520"/>
            </a:xfrm>
          </p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
        from lxml import etree

        for xml in (layout_xml, body_small_xml, body_banner_xml):
            spTree.append(etree.fromstring(xml.strip()))

        # Layout name "Title and Content" 은 _pick_body_layout 의 exact-match
        # 1순위라 이 layout 이 선택됨 — 다른 layouts 제거할 필요 없음.

        tpl_buf = BytesIO()
        tpl.save(tpl_buf)
        raw = tpl_buf.getvalue()

        _set_template_config("pptx")
        with _mock_storage_returns_path(raw):
            content = PPTContent(
                filename="d",
                title="T",
                slides=[PPTSlide(title="Slide Title", bullets=["A", "B", "C"])],
            )
            buf = _build_pptx(content)
        buf.seek(0)
        prs = Presentation(buf)

        body_slide = prs.slides[1]
        # 본문 placeholder (idx=18 banner) 에 bullets 가 들어가지 않아야 한다
        for ph in body_slide.placeholders:
            if ph.placeholder_format.idx == 18:
                assert ph.text_frame.text == "", (
                    "Bullets 가 banner placeholder (idx=18, h=0.8in) 에 잘못 들어감 — "
                    "거버닝 메시지 회귀"
                )

        # bullets 는 textbox shape 에 들어가야 한다
        textboxes_with_content = [
            s
            for s in body_slide.shapes
            if (not s.is_placeholder) and s.has_text_frame and s.text_frame.text
        ]
        assert len(textboxes_with_content) == 1, (
            f"Expected exactly 1 textbox with bullets, got {len(textboxes_with_content)}"
        )
        tb = textboxes_with_content[0]
        assert "A\nB\nC" in tb.text_frame.text

        # textbox 가 widescreen 슬라이드 너비를 거의 다 채워야 한다 (좌측만 차지 회귀 차단)
        from pptx.util import Emu

        slide_w_in = Emu(prs.slide_width).inches
        tb_w_in = Emu(tb.width).inches
        assert tb_w_in > slide_w_in * 0.85, (
            f"Textbox 너비 {tb_w_in:.2f}in 가 슬라이드 너비 {slide_w_in:.2f}in 의 "
            f"85% 미만 — widescreen fallback 회귀"
        )

    def test_pptx_corrupt_template_falls_back_silently(self):
        """A corrupt template should NOT raise — builder falls back to default."""
        _set_template_config("pptx")
        with _mock_storage_returns_path(b"garbage"):
            content = PPTContent(filename="d", title="T", slides=[PPTSlide(title="S")])
            buf = _build_pptx(content)  # must not raise
        buf.seek(0)
        prs = Presentation(buf)
        # Falls back to built-in (Presentation()), so we get 2 slides
        assert len(prs.slides) == 2
