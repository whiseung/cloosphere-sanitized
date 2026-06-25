"""create_pptx 툴 테스트."""

import pytest
from extension_modules.tools.document.pptx_tool import PPTContent, PPTSlide
from pydantic import ValidationError


class TestPPTSchema:
    def test_minimal_valid(self):
        c = PPTContent(
            filename="deck",
            title="My Deck",
            slides=[PPTSlide(title="S1")],
        )
        assert c.slides[0].title == "S1"
        assert c.slides[0].bullets == []
        assert c.slides[0].speaker_notes is None

    def test_empty_slides_rejected(self):
        with pytest.raises(ValidationError):
            PPTContent(filename="d", title="T", slides=[])

    def test_too_many_bullets_rejected(self):
        with pytest.raises(ValidationError):
            PPTSlide(title="S", bullets=["b"] * 21)

    def test_subtitle_optional(self):
        c = PPTContent(filename="d", title="T", slides=[PPTSlide(title="S")])
        assert c.subtitle is None

    def test_speaker_notes_max_length(self):
        with pytest.raises(ValidationError):
            PPTSlide(title="S", speaker_notes="x" * 2001)


from unittest.mock import MagicMock, patch

from extension_modules.tools.document.pptx_tool import (
    PPTX_MIME,
    _build_pptx,
    make_create_pptx,
)
from pptx import Presentation as PptxOpen


class TestPPTBuilder:
    def test_round_trip_title_and_content(self):
        content = PPTContent(
            filename="d",
            title="Quarterly Review",
            subtitle="Q1 2026",
            slides=[
                PPTSlide(title="Intro", bullets=["A", "B", "C"]),
                PPTSlide(title="Data", bullets=["X"], speaker_notes="더 설명"),
            ],
        )
        buf = _build_pptx(content)
        buf.seek(0)
        prs = PptxOpen(buf)

        # 슬라이드 수: title 1 + content 2 = 3
        assert len(prs.slides) == 3
        # 첫 슬라이드 = title slide
        title_slide = prs.slides[0]
        title_text = title_slide.shapes.title.text
        assert title_text == "Quarterly Review"

        # 두 번째 슬라이드 = "Intro"
        s2 = prs.slides[1]
        assert s2.shapes.title.text == "Intro"
        # bullets은 content placeholder의 text frame paragraphs
        body_texts = [
            p.text
            for shape in s2.shapes
            if shape.has_text_frame and shape != s2.shapes.title
            for p in shape.text_frame.paragraphs
        ]
        assert "A" in body_texts and "B" in body_texts and "C" in body_texts

        # 세 번째 슬라이드의 speaker notes
        s3 = prs.slides[2]
        notes = s3.notes_slide.notes_text_frame.text
        assert "더 설명" in notes

    def test_round_trip_minimal(self):
        content = PPTContent(
            filename="d",
            title="T",
            slides=[PPTSlide(title="S")],
        )
        buf = _build_pptx(content)
        buf.seek(0)
        prs = PptxOpen(buf)
        assert len(prs.slides) == 2  # title + 1


class TestMakeCreatePptx:
    def test_returns_structured_tool(self):
        tool = make_create_pptx(user_id="u-1")
        assert tool.name == "create_pptx"
        assert (
            "발표" in tool.description
            or "슬라이드" in tool.description
            or "PPT" in tool.description
        )

    def test_invocation_calls_save(self):
        tool = make_create_pptx(user_id="u-1")
        with patch(
            "extension_modules.tools.document.pptx_tool.save_to_files"
        ) as mock_save:
            mock_save.return_value = MagicMock(
                file_id="f", filename="d.pptx", download_url="/", size_bytes=1
            )
            tool.invoke(
                {
                    "filename": "d",
                    "title": "T",
                    "slides": [{"title": "S"}],
                }
            )
            mock_save.assert_called_once()
            assert mock_save.call_args.kwargs["user_id"] == "u-1"
            assert mock_save.call_args.kwargs["filename"] == "d.pptx"
            assert mock_save.call_args.kwargs["mime"] == PPTX_MIME


# ── Styling tests ─────────────────────────────────────────────────────────────


class TestPPTStyling:
    def _build_and_load(self, content: PPTContent):
        buf = _build_pptx(content)
        buf.seek(0)
        return PptxOpen(buf)

    def test_title_slide_styled(self):
        prs = self._build_and_load(
            PPTContent(
                filename="d",
                title="My Deck",
                subtitle="2026",
                slides=[PPTSlide(title="S")],
            )
        )
        title_slide = prs.slides[0]
        # 첫 run 이 deep blue + bold + 40pt
        title_run = title_slide.shapes.title.text_frame.paragraphs[0].runs[0]
        assert title_run.font.bold is True
        assert title_run.font.size.pt == 40
        # Subtitle: teal + 24pt
        subtitle_run = title_slide.placeholders[1].text_frame.paragraphs[0].runs[0]
        assert subtitle_run.font.size.pt == 24

    def test_content_slide_heading_and_bullets_styled(self):
        prs = self._build_and_load(
            PPTContent(
                filename="d",
                title="T",
                slides=[PPTSlide(title="Section", bullets=["item1", "item2"])],
            )
        )
        content_slide = prs.slides[1]
        # Heading: bold + 28pt
        heading_run = content_slide.shapes.title.text_frame.paragraphs[0].runs[0]
        assert heading_run.font.bold is True
        assert heading_run.font.size.pt == 28
        # Bullet: 18pt (모든 paragraph)
        body_tf = None
        for shape in content_slide.shapes:
            if shape.has_text_frame and shape != content_slide.shapes.title:
                body_tf = shape.text_frame
                break
        assert body_tf is not None
        for paragraph in body_tf.paragraphs:
            for run in paragraph.runs:
                assert run.font.size.pt == 18

    def test_table_slide_renders_with_header_styling(self):
        prs = self._build_and_load(
            PPTContent(
                filename="d",
                title="T",
                slides=[
                    PPTSlide(
                        title="비교",
                        table=[["구분", "A", "B"], ["가격", "100", "200"]],
                    )
                ],
            )
        )
        # 슬라이드 1 = title, 슬라이드 2 = 표 슬라이드
        table_slide = prs.slides[1]
        tables = [s for s in table_slide.shapes if s.has_table]
        assert len(tables) == 1, "슬라이드에 표가 1개 추가되어야 함"
        tbl = tables[0].table
        assert len(tbl.rows) == 2
        assert len(tbl.columns) == 3
        assert tbl.cell(0, 0).text == "구분"
        assert tbl.cell(1, 1).text == "100"

        # 헤더 행: bold + 흰 글자
        header_run = tbl.cell(0, 0).text_frame.paragraphs[0].runs[0]
        assert header_run.font.bold is True
        # 데이터 행: bold 아님
        data_run = tbl.cell(1, 0).text_frame.paragraphs[0].runs[0]
        assert data_run.font.bold is not True


class TestPPTSlideContentChoice:
    def test_bullets_dropped_when_table_set(self):
        # bullets 와 table 동시 지정 시 table 우선, bullets 자동 비움
        s = PPTSlide(title="S", bullets=["b1", "b2"], table=[["A"], ["1"]])
        assert s.bullets == []
        assert s.table == [["A"], ["1"]]

    def test_table_row_cap_enforced(self):
        # max_length=30 — 31행 거부
        with pytest.raises(ValidationError):
            PPTSlide(title="S", table=[["H"]] + [["x"]] * 31)
